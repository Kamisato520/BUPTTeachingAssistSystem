import logging
from PIL import Image
import io
import os
from bs4 import BeautifulSoup
import markdown

logger = logging.getLogger(__name__)

class FileProcessor:
    ALLOWED_EXTENSIONS = {
        'image': {'png', 'jpg', 'jpeg', 'gif', 'bmp'},
        'document': {'pdf', 'docx', 'doc', 'pptx', 'ppt'},
        'text': {'txt', 'md', 'html'}
    }

    @staticmethod
    def is_allowed_file(filename):
        return '.' in filename and \
            filename.rsplit('.', 1)[1].lower() in \
            {ext for exts in FileProcessor.ALLOWED_EXTENSIONS.values() for ext in exts}

    @staticmethod
    def extract_text_from_file(file, file_type):
        """从不同类型的文件中提取文本"""
        try:
            if file_type in FileProcessor.ALLOWED_EXTENSIONS['image']:
                return FileProcessor._process_image(file)
            elif file_type == 'pdf':
                return FileProcessor._process_pdf(file)
            elif file_type in ['docx', 'doc']:
                return FileProcessor._process_word(file)
            elif file_type in ['pptx', 'ppt']:
                return FileProcessor._process_powerpoint(file)
            elif file_type == 'md':
                return FileProcessor._process_markdown(file)
            elif file_type == 'html':
                return FileProcessor._process_html(file)
            elif file_type == 'txt':
                return file.read().decode('utf-8')
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
        except Exception as e:
            logger.error(f"文件处理错误: {str(e)}")
            return f"文件处理错误: {str(e)}"

    @staticmethod
    def _process_image(file):
        """处理图片文件"""
        try:
            # 尝试导入 pytesseract
            try:
                import pytesseract
                image = Image.open(io.BytesIO(file.read()))
                text = pytesseract.image_to_string(image, lang='chi_sim+eng')
                return text
            except ImportError:
                return "[图片内容] - OCR组件未安装，无法提取文字"
        except Exception as e:
            logger.error(f"图片处理错误: {str(e)}")
            return f"[图片内容] - 处理错误: {str(e)}"

    @staticmethod
    def _process_pdf(file):
        """处理PDF文件"""
        try:
            # 尝试导入 PyMuPDF
            try:
                import fitz
                pdf_document = fitz.open(stream=file.read(), filetype="pdf")
                text = ""
                for page in pdf_document:
                    text += page.get_text()
                return text
            except ImportError:
                return "[PDF内容] - PDF处理组件未安装，无法提取文字"
        except Exception as e:
            logger.error(f"PDF处理错误: {str(e)}")
            return f"[PDF内容] - 处理错误: {str(e)}"

    @staticmethod
    def _process_word(file):
        """处理Word文件"""
        try:
            # 尝试导入 python-docx
            try:
                from docx import Document
                doc = Document(io.BytesIO(file.read()))
                return "\n".join([paragraph.text for paragraph in doc.paragraphs])
            except ImportError:
                return "[Word文档内容] - Word处理组件未安装，无法提取文字"
        except Exception as e:
            logger.error(f"Word处理错误: {str(e)}")
            return f"[Word文档内容] - 处理错误: {str(e)}"

    @staticmethod
    def _process_powerpoint(file):
        """处理PowerPoint文件"""
        try:
            # 尝试导入 python-pptx
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file.read()))
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            except ImportError:
                return "[PPT内容] - PPT处理组件未安装，无法提取文字"
        except Exception as e:
            logger.error(f"PowerPoint处理错误: {str(e)}")
            return f"[PPT内容] - 处理错误: {str(e)}"

    @staticmethod
    def _process_markdown(file):
        """处理Markdown文件"""
        try:
            md_text = file.read().decode('utf-8')
            html = markdown.markdown(md_text)
            soup = BeautifulSoup(html, 'html.parser')
            return soup.get_text()
        except Exception as e:
            logger.error(f"Markdown处理错误: {str(e)}")
            return f"[Markdown内容] - 处理错误: {str(e)}"

    @staticmethod
    def _process_html(file):
        """处理HTML文件"""
        try:
            html_content = file.read().decode('utf-8')
            soup = BeautifulSoup(html_content, 'html.parser')
            return soup.get_text()
        except Exception as e:
            logger.error(f"HTML处理错误: {str(e)}")
            return f"[HTML内容] - 处理错误: {str(e)}" 